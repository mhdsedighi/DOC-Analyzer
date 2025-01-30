import os
import json
import tkinter as tk
from tkinter import filedialog, scrolledtext, messagebox, ttk
import PyPDF2
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import ollama
from pytesseract import image_to_string
from pdf2image import convert_from_path
from PIL import Image

# File paths for saving user data and cache
USER_DATA_FILE = "user_data.json"
DOCUMENT_CACHE_FILE = "document_cache.json"

# Supported file extensions
SUPPORTED_EXTENSIONS = [".pdf", ".docx", ".txt", ".xlsx", ".pptx"]

chat_history_list = []  # List to store the conversation history

# Load user data (last folder path, last selected model, and temperature)
def load_user_data():
    if os.path.exists(USER_DATA_FILE):
        with open(USER_DATA_FILE, "r") as file:
            return json.load(file)
    return {"last_folder": "", "last_model": "", "temperature": 0.7}  # Default temperature

# Save user data (last folder path, last selected model, and temperature)
def save_user_data(last_folder, last_model, temperature):
    user_data = {
        "last_folder": last_folder,
        "last_model": last_model,
        "temperature": temperature
    }
    with open(USER_DATA_FILE, "w") as file:
        json.dump(user_data, file, indent=4)

# Load the document cache
def load_document_cache():
    if os.path.exists(DOCUMENT_CACHE_FILE):
        with open(DOCUMENT_CACHE_FILE, "r") as file:
            return json.load(file)
    return {}

# Save the document cache
def save_document_cache(cache):
    with open(DOCUMENT_CACHE_FILE, "w") as file:
        json.dump(cache, file, indent=4)

# Function to fetch installed Ollama models
def fetch_installed_models():
    try:
        models = ollama.list()  # Fetch the list of installed models
        return [model["name"] for model in models.get("models", [])]
    except Exception as e:
        messagebox.showerror("Error", f"Failed to fetch models: {e}")
        return []

# Function to extract text from a PDF file using PyPDF2 and OCR
def extract_text_from_pdf(pdf_path):
    text = ""
    unreadable_pages = 0
    total_pages = 0

    try:
        # Try extracting text using PyPDF2
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            total_pages = len(reader.pages)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
                else:
                    unreadable_pages += 1
    except Exception as e:
        print(f"PyPDF2 extraction failed: {e}")

    # If PyPDF2 fails or some pages are unreadable, use OCR
    if unreadable_pages > 0:
        try:
            # Convert PDF pages to images
            images = convert_from_path(pdf_path)
            for image in images:
                ocr_text = image_to_string(image)
                text += ocr_text + "\n"
        except Exception as e:
            print(f"OCR extraction failed: {e}")

    # Calculate the percentage of readable content
    readable_percentage = 100 - (unreadable_pages / total_pages * 100) if total_pages > 0 else 100

    # Count the number of words
    word_count = len(text.split())

    return text, word_count, readable_percentage

# Function to extract text from a DOCX file
def extract_text_from_docx(docx_path):
    try:
        doc = Document(docx_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        word_count = len(text.split())
        return text, word_count, 100  # Assume 100% readable for DOCX
    except Exception as e:
        print(f"DOCX extraction failed: {e}")
        return "", 0, 0  # Assume 0% readable if extraction fails

# Function to extract text from a TXT file
def extract_text_from_txt(txt_path):
    try:
        with open(txt_path, "r", encoding="utf-8") as file:
            text = file.read()
        word_count = len(text.split())
        return text, word_count, 100  # Assume 100% readable for TXT
    except Exception as e:
        print(f"TXT extraction failed: {e}")
        return "", 0, 0  # Assume 0% readable if extraction fails

# Function to extract text from an XLSX file
def extract_text_from_xlsx(xlsx_path):
    try:
        workbook = load_workbook(xlsx_path)
        text = ""
        for sheet in workbook:
            for row in sheet.iter_rows(values_only=True):
                text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
        word_count = len(text.split())
        return text, word_count, 100  # Assume 100% readable for XLSX
    except Exception as e:
        print(f"XLSX extraction failed: {e}")
        return "", 0, 0  # Assume 0% readable if extraction fails

# Function to extract text from a PPTX file
def extract_text_from_pptx(pptx_path):
    try:
        presentation = Presentation(pptx_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        word_count = len(text.split())
        return text, word_count, 100  # Assume 100% readable for PPTX
    except Exception as e:
        print(f"PPTX extraction failed: {e}")
        return "", 0, 0  # Assume 0% readable if extraction fails

# Function to extract text from a document based on its file type
def extract_text_from_document(file_path):
    if file_path.endswith(".pdf"):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith(".docx"):
        return extract_text_from_docx(file_path)
    elif file_path.endswith(".txt"):
        return extract_text_from_txt(file_path)
    elif file_path.endswith(".xlsx"):
        return extract_text_from_xlsx(file_path)
    elif file_path.endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    else:
        return "", 0, 0  # Unsupported file type

# Function to read all documents in a folder
def read_documents(folder_path):
    cache = load_document_cache()
    all_text = ""
    new_files_read = 0  # Track the number of new files read

    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        file_ext = os.path.splitext(filename)[1].lower()

        if file_ext in SUPPORTED_EXTENSIONS:
            last_modified = os.path.getmtime(file_path)
            
            # Check if the file is in the cache and hasn't been modified
            if file_path in cache and cache[file_path]["last_modified"] == last_modified:
                text = cache[file_path]["text"]
                word_count = cache[file_path]["word_count"]
                readable_percentage = cache[file_path]["readable_percentage"]
            else:
                # Extract text and update the cache
                text, word_count, readable_percentage = extract_text_from_document(file_path)
                cache[file_path] = {
                    "last_modified": last_modified,
                    "text": text,
                    "word_count": word_count,
                    "readable_percentage": readable_percentage
                }
                new_files_read += 1  # Increment the counter for new files
            
            all_text += f"--- {filename} ---\n{text}\n\n"
            chat_history.config(state=tk.NORMAL)
            chat_history.insert(tk.END, f"Looked at: {filename}\n")
            chat_history.insert(tk.END, f"Word count: {word_count}\n")
            chat_history.insert(tk.END, f"Readable content: {readable_percentage:.2f}%\n\n")
            chat_history.config(state=tk.DISABLED)
    
    # Save the updated cache
    save_document_cache(cache)
    return all_text, new_files_read

# Function to handle the chat with the AI
def chat_with_ai():
    user_input = user_input_box.get("1.0", tk.END).strip()
    if user_input:
        chat_history.config(state=tk.NORMAL)
        chat_history.insert(tk.END, f"You: {user_input}\n")
        
        # Add the user's message to the chat history list
        chat_history_list.append({"role": "user", "content": user_input})
        
        # Combine the document text with the chat history
        full_prompt = f"{document_text}\n\n" + "\n".join([f"{msg['role']}: {msg['content']}" for msg in chat_history_list])
        
        try:
            # Send the prompt to the AI using the selected model
            selected_model = model_var.get()
            response = ollama.chat(
                model=selected_model,
                messages=[{"role": "user", "content": full_prompt}],
                options={"temperature": temperature_scale.get()}  # Use the slider value
            )
            
            ai_response = response['message']['content']
            chat_history.insert(tk.END, f"AI: {ai_response}\n")
            
            # Add the AI's response to the chat history list
            chat_history_list.append({"role": "assistant", "content": ai_response})
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred: {e}")
        
        chat_history.config(state=tk.DISABLED)
        user_input_box.delete("1.0", tk.END)
        
        # Save the last used model, folder path, and temperature
        save_user_data(folder_path_entry.get().strip(), selected_model, temperature_scale.get())

# Function to clear the chat history box
def clear_chat_history():
    chat_history.config(state=tk.NORMAL)
    chat_history.delete("1.0", tk.END)
    chat_history.config(state=tk.DISABLED)
    chat_history_list.clear()  # Clear the chat history list

# Function to set the folder path
def set_folder_path():
    folder_path = folder_path_entry.get().strip()
    if not folder_path:
        messagebox.showwarning("Warning", "Please enter a folder path.")
        return
    
    if not os.path.isdir(folder_path):
        messagebox.showerror("Error", "Invalid folder path.")
        return
    
    global document_text
    document_text, new_files_read = read_documents(folder_path)
    chat_history.config(state=tk.NORMAL)
    chat_history.insert(tk.END, f"Documents reading finished. {new_files_read} new files were processed.\n")
    chat_history.insert(tk.END, "You can now chat with the AI.\n")
    chat_history.config(state=tk.DISABLED)
    
    # Save the folder path, last used model, and temperature
    save_user_data(folder_path, model_var.get(), temperature_scale.get())

# Create the main window
root = tk.Tk()
root.title("AI Document Analyzer")

# Set the background color of the main window to gray
root.configure(bg="#333333")

# Configure grid weights to make the chat history box resizable
root.grid_rowconfigure(1, weight=1)  # Make row 1 (chat history) resizable
root.grid_columnconfigure(0, weight=1)  # Make column 0 resizable

# Fetch installed Ollama models
installed_models = fetch_installed_models()

# Load user data (last folder path, last selected model, and temperature)
user_data = load_user_data()
last_folder = user_data.get("last_folder", "")
last_model = user_data.get("last_model", "")
last_temperature = user_data.get("temperature", 0.7)  # Default temperature

# Dropdown for model selection
model_var = tk.StringVar(root)
model_var.set(last_model if last_model in installed_models else (installed_models[0] if installed_models else "No models found"))
model_dropdown = ttk.Combobox(root, textvariable=model_var, values=installed_models, state="readonly")
model_dropdown.grid(row=0, column=3, padx=10, pady=10, sticky="e")

# Folder path entry
folder_path_entry = ttk.Entry(root, width=50)
folder_path_entry.grid(row=0, column=0, padx=10, pady=10, sticky="ew")

# Set the last used folder path
folder_path_entry.insert(0, last_folder)

# Browse button
browse_button = ttk.Button(root, text="Browse", command=lambda: folder_path_entry.insert(0, filedialog.askdirectory()))
browse_button.grid(row=0, column=1, padx=10, pady=10)

# Read button
read_button = ttk.Button(root, text="Read Documents", command=set_folder_path)
read_button.grid(row=0, column=2, padx=10, pady=10)

# Chat history display
chat_history = scrolledtext.ScrolledText(root, width=80, height=20, state=tk.DISABLED, bg="#444444", fg="white", insertbackground="white")
chat_history.grid(row=1, column=0, columnspan=4, padx=10, pady=10, sticky="nsew")

# User input box
user_input_box = tk.Text(root, width=60, height=3, bg="#444444", fg="white", insertbackground="white")
user_input_box.grid(row=2, column=0, padx=10, pady=10, sticky="ew")

# Send button
send_button = ttk.Button(root, text="Ask AI", command=chat_with_ai)
send_button.grid(row=2, column=1, padx=10, pady=10)

# Clear button
clear_button = ttk.Button(root, text="Clear", command=clear_chat_history)
clear_button.grid(row=2, column=2, padx=10, pady=10)

# Temperature slider
temperature_label = ttk.Label(root, text="Temperature:", background="#333333", foreground="white")
temperature_label.grid(row=2, column=3, padx=10, pady=10)

temperature_scale = tk.Scale(root, from_=0.0, to=1.0, resolution=0.1, orient=tk.HORIZONTAL, bg="#333333", fg="white", troughcolor="#444444")
temperature_scale.set(last_temperature)  # Set the last used temperature
temperature_scale.grid(row=2, column=4, padx=10, pady=10)

# Run the application
root.mainloop()