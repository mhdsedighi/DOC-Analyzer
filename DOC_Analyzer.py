import os
import json
import tkinter as tk
from tkinter import filedialog, scrolledtext, messagebox, ttk
from tkinter import Menu
import enchant  # For spell checking
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import ollama
from pytesseract import image_to_string
from pdf2image import convert_from_path
import win32com.client  # For handling old Microsoft Office formats
from io import BytesIO
import fitz  # PyMuPDF
import pdfplumber
import base64
from io import BytesIO
from PIL import Image

# Define the cache folder and ensure it exists
if not os.path.exists("cache"):
    os.makedirs("cache")

# File paths for saving user data and cache
USER_DATA_FILE = os.path.join("cache", "user_data.json")
DOCUMENT_CACHE_FILE = os.path.join("cache", "document_cache.json")

# Supported file extensions
SUPPORTED_EXTENSIONS = [".pdf", ".docx", ".doc", ".txt", ".xlsx", ".xls", ".pptx", ".ppt"]

# List of models that support image processing
IMAGE_SUPPORTED_MODELS = ["llava", "bakllava", "cogvlm"]  # Add other image-supported models here

chat_history_list = []  # List to store the conversation history

# Load user data (last folder path, last selected model, and temperature)
def load_user_data():
    if os.path.exists(USER_DATA_FILE):
        with open(USER_DATA_FILE, "r") as file:
            data = json.load(file)
            # Ensure the "last_folders" and "do_mention_page" keys exist in the loaded data
            if "last_folders" not in data:
                data["last_folders"] = []
            if "do_mention_page" not in data:
                data["do_mention_page"] = False  # Default value if not present
            return data
    return {
        "last_folder": "",
        "last_folders": [],
        "last_model": "",
        "temperature": 0.7,
        "do_mention_page": False  # Default value for the checkbox
    }

# Save user data (last folder path, last selected model, temperature, and last 10 folders)
def save_user_data(last_folder=None, last_model=None, temperature=None, last_folders=None, do_mention_page=None):
    user_data = load_user_data()
    if last_folder is not None:
        user_data["last_folder"] = last_folder
    if last_model is not None:
        user_data["last_model"] = last_model
    if temperature is not None:
        user_data["temperature"] = temperature
    if last_folders is not None:
        user_data["last_folders"] = last_folders
    if do_mention_page is not None:
        user_data["do_mention_page"] = do_mention_page

    # Update the last 10 folders list if a new folder is added
    if last_folder and last_folder not in user_data.get("last_folders", []):
        if "last_folders" not in user_data:
            user_data["last_folders"] = []
        user_data["last_folders"].insert(0, last_folder)
        user_data["last_folders"] = user_data["last_folders"][:10]  # Keep only the last 10

    with open(USER_DATA_FILE, "w") as file:
        json.dump(user_data, file, indent=4)

    # Update the last 10 folders list
    if last_folder:
        if last_folder in user_data["last_folders"]:
            user_data["last_folders"].remove(last_folder)  # Remove if already exists
        user_data["last_folders"].insert(0, last_folder)  # Add to the beginning
        user_data["last_folders"] = user_data["last_folders"][:10]  # Keep only the last 10

    with open(USER_DATA_FILE, "w") as file:
        json.dump(user_data, file, indent=4)

# Load the document cache
def load_document_cache():
    if os.path.exists(DOCUMENT_CACHE_FILE):
        with open(DOCUMENT_CACHE_FILE, "r") as file:
            return json.load(file)
    return {}

# Save the document cache
def save_document_cache(cache):
    with open(DOCUMENT_CACHE_FILE, "w") as file:
        json.dump(cache, file, indent=4)

# Function to fetch installed Ollama models
def fetch_installed_models():
    try:
        models_response = ollama.list()  # Fetch the list of installed models
        # print("DEBUG: Response from ollama.list():", models_response)  # Debugging output
        
        models_list = models_response.models  # Access the 'models' attribute directly

        if not isinstance(models_list, list):
            raise ValueError("Unexpected response format: 'models' key is not a list")

        # Extract model names
        model_names = [getattr(model, "model", "Unknown") for model in models_list]

        return model_names

    except Exception as e:
        messagebox.showerror("Error", f"Failed to fetch models: {e}")
        return []

# Function to handle folder browsing
def browse_folder():
    folder_path = filedialog.askdirectory()
    if folder_path:
        folder_path_entry.delete(0, tk.END)  # Clear the entry widget
        folder_path_entry.insert(0, folder_path)  # Insert the new folder path
        update_folder_dropdown()  # Update the dropdown with the new path
        save_user_data(last_folder=folder_path)  #save the new folder path


# Function to extract text from a DOCX file
def extract_content_from_docx(docx_path):
    text_content = []  # Stores extracted text
    image_content = []  # Stores extracted images (DOCX may contain embedded images)
    word_count = 0

    try:
        doc = Document(docx_path)
        page_num = 1  # Since DOCX doesn’t have pages, treat all content as page 1

        # Extract text
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        if text:
            text_content.append({"page": page_num, "content": text})
            word_count = len(text.split())

        # Extract images if enabled
        if do_read_image:
            for rel in doc.part.rels:
                if "image" in doc.part.rels[rel].target_ref:
                    image_data = doc.part.rels[rel].target_part.blob
                    img_base64 = base64.b64encode(image_data).decode("utf-8")
                    image_content.append({"page": page_num, "content": img_base64})

        return text_content, image_content, word_count, 100  # Assume 100% readability

    except Exception as e:
        print(f"DOCX extraction failed: {e}")
        return [], [], 0, 0  # Assume 0% readable if extraction fails


# Function to extract text and images from a DOC file (old Word format)
def extract_content_from_doc(doc_path):
    text_content = []  # Stores extracted text
    image_content = []  # Stores extracted images
    word_count = 0

    try:
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(doc_path)

        # Extract text
        text = doc.Content.Text.strip()
        page_num = 1  # DOC files don’t have clear pages, so all content is assigned to page 1

        if text:
            text_content.append({"page": page_num, "content": text})
            word_count = len(text.split())

        # Extract images if enabled
        if do_read_image:
            for shape in doc.InlineShapes:
                if shape.Type == 3:  # Type 3 represents images
                    image = shape.PictureFormat
                    image_data = image.SaveAsPicture()  # Save the image temporarily
                    with open(image_data, "rb") as img_file:
                        img_base64 = base64.b64encode(img_file.read()).decode("utf-8")
                        image_content.append({"page": page_num, "content": img_base64})

        doc.Close(False)
        word.Quit()

        return text_content, image_content, word_count, 100  # Assume 100% readable for DOC files

    except Exception as e:
        print(f"DOC extraction failed: {e}")
        return [], [], 0, 0  # Assume 0% readable if extraction fails

# Function to extract text from a TXT file
def extract_content_from_txt(txt_path):
    text_content = []  # Stores extracted text
    image_content = []  # No images for TXT files
    word_count = 0

    try:
        with open(txt_path, "r", encoding="utf-8") as file:
            text = file.read().strip()

        page_num = 1  # Assume the entire document is on one page
        if text:
            text_content.append({"page": page_num, "content": text})
            word_count = len(text.split())

        return text_content, image_content, word_count, 100  # Assume 100% readability

    except Exception as e:
        print(f"TXT extraction failed: {e}")
        return [], [], 0, 0  # Assume 0% readable if extraction fails


# Function to extract text from an XLSX file
def extract_content_from_xlsx(xlsx_path):
    text_content = []  # Stores extracted text per sheet (as pages)
    image_content = []  # No images for XLSX files (unless explicitly handled)
    word_count = 0

    try:
        workbook = load_workbook(xlsx_path)

        for page_num, sheet in enumerate(workbook.sheetnames, start=1):
            sheet_text = ""

            # Read each row in the sheet
            for row in workbook[sheet].iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    sheet_text += row_text + "\n"

            if sheet_text.strip():
                text_content.append({"page": page_num, "content": sheet_text})
                word_count += len(sheet_text.split())

        return text_content, image_content, word_count, 100  # Assume 100% readability

    except Exception as e:
        print(f"XLSX extraction failed: {e}")
        return [], [], 0, 0  # Assume 0% readable if extraction fails


# Function to extract text from an XLS file (old Excel format)
def extract_content_from_xls(xls_path):
    text_content = []  # Stores extracted text per sheet (as pages)
    image_content = []  # No images for XLS files
    word_count = 0

    try:
        excel = win32com.client.Dispatch("Excel.Application")
        excel.Visible = False
        workbook = excel.Workbooks.Open(xls_path)

        for page_num, sheet in enumerate(workbook.Sheets, start=1):
            sheet_text = ""

            # Read each row in the sheet
            for row in sheet.UsedRange.Rows:
                row_text = " ".join([str(cell) for cell in row.Value if cell is not None])
                if row_text.strip():
                    sheet_text += row_text + "\n"

            if sheet_text.strip():
                text_content.append({"page": page_num, "content": sheet_text})
                word_count += len(sheet_text.split())

        workbook.Close(False)  # Close without saving
        excel.Quit()

        return text_content, image_content, word_count, 100  # Assume 100% readability

    except Exception as e:
        print(f"XLS extraction failed: {e}")
        return [], [], 0, 0  # Assume 0% readable if extraction fails



# Function to extract text and images from a PPTX file
def extract_content_from_pptx(pptx_path):
    text_content = []  # Stores extracted text
    image_content = []  # Stores extracted images
    word_count = 0

    try:
        presentation = Presentation(pptx_path)

        for slide_num, slide in enumerate(presentation.slides, start=1):
            slide_text = ""

            # Extract text from each shape in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"

            if slide_text.strip():
                text_content.append({"page": slide_num, "content": slide_text})
                word_count += len(slide_text.split())

            # Extract images if enabled
            if do_read_image:
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        image = shape.image
                        image_stream = BytesIO(image.blob)
                        img_base64 = base64.b64encode(image_stream.getvalue()).decode("utf-8")
                        image_content.append({"page": slide_num, "content": img_base64})

        return text_content, image_content, word_count, 100  # Assume 100% readable for PPTX

    except Exception as e:
        print(f"PPTX extraction failed: {e}")
        return [], [], 0, 0  # Assume 0% readable if extraction fails

# Function to extract text and images from a PPT file (old PowerPoint format)
def extract_content_from_ppt(ppt_path):
    text_content = []  # Stores extracted text
    image_content = []  # Stores extracted images
    word_count = 0

    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = False
        presentation = powerpoint.Presentations.Open(ppt_path, WithWindow=False)

        for slide_num, slide in enumerate(presentation.Slides, start=1):
            slide_text = ""

            # Extract text from each shape in the slide
            for shape in slide.Shapes:
                if hasattr(shape, "TextFrame") and shape.TextFrame.HasText:
                    slide_text += shape.TextFrame.TextRange.Text + "\n"

            if slide_text.strip():
                text_content.append({"page": slide_num, "content": slide_text})
                word_count += len(slide_text.split())

            # Extract images if enabled
            if do_read_image:
                for shape in slide.Shapes:
                    if shape.Type == 13:  # Type 13 represents pictures
                        image = shape.PictureFormat
                        temp_path = os.path.join(tempfile.gettempdir(), f"ppt_image_{slide_num}.png")
                        image.SaveAsFile(temp_path)  # Save the image temporarily
                        with open(temp_path, "rb") as img_file:
                            img_base64 = base64.b64encode(img_file.read()).decode("utf-8")
                            image_content.append({"page": slide_num, "content": img_base64})
                        os.remove(temp_path)  # Clean up the temporary file

        presentation.Close()
        powerpoint.Quit()

        return text_content, image_content, word_count, 100  # Assume 100% readable for PPT

    except Exception as e:
        print(f"PPT extraction failed: {e}")
        return [], [], 0, 0  # Assume 0% readable if extraction fails


# Function to extract text from a document based on its file type
def extract_content_from_document(file_path):
    if file_path.endswith(".pdf"):
        return extract_content_from_pdf(file_path)
    elif file_path.endswith(".docx"):
        return extract_content_from_docx(file_path)
    elif file_path.endswith(".doc"):
        return extract_content_from_doc(file_path)
    elif file_path.endswith(".txt"):
        return extract_content_from_txt(file_path)
    elif file_path.endswith(".xlsx"):
        return extract_content_from_xlsx(file_path)
    elif file_path.endswith(".xls"):
        return extract_content_from_xls(file_path)
    elif file_path.endswith(".pptx"):
        return extract_content_from_pptx(file_path)
    elif file_path.endswith(".ppt"):
        return extract_content_from_ppt(file_path)
    else:
        return [],[], 0, 0  # Unsupported file type


# Function to extract text and images from a PDF file
def extract_content_from_pdf(pdf_path):
    text_content = []  # Stores extracted text per page
    image_content = []  # Stores extracted images per page if enabled
    unreadable_pages = 0
    total_pages = 0
    word_count = 0

    try:
        # Open the PDF file
        with pdfplumber.open(pdf_path) as pdf:
            total_pages = len(pdf.pages)
            doc = fitz.open(pdf_path)  # Open with PyMuPDF for image extraction

            for page_num in range(total_pages):
                page_text = pdf.pages[page_num].extract_text()

                if page_text:
                    text_content.append({"page": page_num + 1, "content": page_text})
                    word_count += len(page_text.split())
                else:
                    unreadable_pages += 1

                if do_read_image:
                    # Extract raster images from the page
                    fitz_page = doc.load_page(page_num)
                    image_list = fitz_page.get_images(full=True)

                    for img_index, img in enumerate(image_list):
                        xref = img[0]  # XREF of the image
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]

                        # Convert image bytes to base64
                        image = Image.open(BytesIO(image_bytes))
                        buffered = BytesIO()
                        image.save(buffered, format="PNG")
                        img_base64 = base64.b64encode(buffered.getvalue()).decode("utf-8")

                        image_content.append({"page": page_num + 1, "content": img_base64})

        readable_percentage = 100 - (unreadable_pages / total_pages * 100) if total_pages > 0 else 100
        return text_content, image_content, word_count, readable_percentage

    except Exception as e:
        print(f"Extraction failed: {e}")
        return [], [], 0, 0


# Function to read all documents in a folder
def read_documents(folder_path):
    document_text = '' #this appers to fix adding extra texts when changing folder
    cache = load_document_cache()
    new_files_read = 0  # Track the number of new files read
    document_images = []  # List to store images from documents

    # Set the introductory line based on the checkbox state
    if do_mention_var.get():
        all_text = """Below are the contents of several files for analysis.
                    The filename and page number are mentioned with each content block.
                    When responding, reference the source document and page number.
                    Example: 'The data shows an increase in sales [report.pdf, page 3]'."""
    else:
        all_text = "Below are the extracted contents of the documents:\n\n"

    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        file_ext = os.path.splitext(filename)[1].lower()

        if file_ext in SUPPORTED_EXTENSIONS:
            last_modified = os.path.getmtime(file_path)

            # Check if the file is in the cache and hasn't been modified
            if file_path in cache and cache[file_path]["last_modified"] == last_modified:
                text_content = cache[file_path]["text_content"]
                image_content = cache[file_path]["image_content"] if do_read_image else []
                word_count = cache[file_path]["word_count"]
                readable_percentage = cache[file_path]["readable_percentage"]
            else:
                text_content, image_content, word_count, readable_percentage = extract_content_from_document(file_path)
                cache[file_path] = {
                    "last_modified": last_modified,
                    "text_content": text_content,
                    "image_content": image_content if do_read_image else [],
                    "word_count": word_count,
                    "readable_percentage": readable_percentage
                }
                new_files_read += 1

            all_text += f"--- Document: {filename} ---\n"

            for item in text_content:
                all_text += f"--- Page {item['page']} ---\n{item['content']}\n\n"

            if do_read_image:
                for item in image_content:
                    all_text += f"[Image on Page {item['page']}]\n"
                    document_images.append(item)

            chat_history.config(state=tk.NORMAL)
            chat_history.insert(tk.END, f"Looked at: {filename}\n", "fileread_tag")
            chat_history.insert(tk.END, f"Word count: {word_count}\n", "fileread_tag")
            chat_history.insert(tk.END, f"Readable content: {readable_percentage:.2f}%\n\n", "fileread_tag")
            chat_history.config(state=tk.DISABLED)

    save_document_cache(cache)
    return all_text, new_files_read, document_images


# Function to handle the chat with the AI
def chat_with_ai():
    user_input = user_input_box.get("1.0", tk.END).strip()
    global previous_message
    global do_revise
    previous_message=user_input

    if do_revise:  #removing what has been revised from prompt
        if chat_history_list:
            chat_history_list.pop()
            chat_history_list.pop()
        do_revise=False

    if user_input:
        chat_history.config(state=tk.NORMAL)
        
        # Insert the user's question
        chat_history.insert(tk.END, "You: ", "user_tag")  # Tag for user text
        chat_history.insert(tk.END, f"{user_input}\n", "user_question")  # Tag for user question
        
        # Add the user's message to the chat history list
        chat_history_list.append({"role": "user", "content": user_input})
        
        # Combine the document text with the chat history
        full_prompt = f"{document_text}\n\n" + "\n".join([f"{msg['role']}: {msg['content']}" for msg in chat_history_list])
        
        try:
            # Send the prompt to the AI using the selected model
            selected_model = model_var.get()
            messages = [{"role": "user", "content": full_prompt}]

            # If the model supports images, include them in the messages
            if selected_model in IMAGE_SUPPORTED_MODELS and document_images:
                for img_base64 in document_images:
                    messages.append({"role": "user", "content": f"data:image/png;base64,{img_base64}"})

            response = ollama.chat(
                model=selected_model,
                messages=messages,
                options={"temperature": temperature_scale.get()}  # Use the slider value
            )
            
            ai_response = response['message']['content']
            chat_history.insert(tk.END, "AI: ", "ai_tag")  # Tag for "AI:" 
            chat_history.insert(tk.END, f"{ai_response}\n", "ai_response")  # Tag for AI response
            
            # Add a horizontal line after the AI's response
            chat_history.insert(tk.END, "---\n", "separator")  # Tag for the separator line
            
            # Add the AI's response to the chat history list
            chat_history_list.append({"role": "assistant", "content": ai_response})
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred: {e}")
        
        chat_history.config(state=tk.DISABLED)
        user_input_box.delete("1.0", tk.END)

        chat_history.see(tk.END) #scrolling down
        
        # Save the last used model, folder path, temperature, and checkbox state
        save_user_data(
            folder_path_entry.get().strip(),
            selected_model,
            temperature_scale.get(),
            do_mention_page=do_mention_var.get()
        )

# Function to clear the chat history box
def clear_chat_history():
    chat_history.config(state=tk.NORMAL)
    chat_history.delete("1.0", tk.END)
    chat_history.config(state=tk.DISABLED)
    chat_history_list.clear()  # Clear the chat history list

# Function to set the folder path
def set_folder_path():
    global document_text
    document_text = ""  # Clear the previous document text
    global document_images
    document_images= []

    folder_path = folder_path_entry.get().strip()
    if not folder_path:
        messagebox.showwarning("Warning", "Please enter a folder path.")
        return
    
    if not os.path.isdir(folder_path):
        messagebox.showerror("Error", "Invalid folder path.")
        return
    
    # Read documents from the new folder
    document_text, new_files_read, document_images =  read_documents(folder_path)
    chat_history.config(state=tk.NORMAL)
    chat_history.insert(tk.END, f"Documents reading finished. {new_files_read} new files were processed.\n")
    chat_history.insert(tk.END, "You can now chat with the A.I.\n")
    chat_history.config(state=tk.DISABLED)
    chat_history.see(tk.END) #scrolling down
    
    # Save the folder path, last used model, and temperature
    save_user_data(folder_path, model_var.get(), temperature_scale.get())

    # Update the dropdown with the latest folder paths
    update_folder_dropdown()

# Function to update the folder dropdown with the last 10 used paths
def update_folder_dropdown():
    user_data = load_user_data()
    folder_path_dropdown["values"] = user_data.get("last_folders", [])
    folder_path_dropdown.set(folder_path_entry.get().strip())  # Sync dropdown with entry

# Function to handle folder path selection from the dropdown
def on_folder_select(event):
    selected_path = folder_path_dropdown.get()
    folder_path_entry.delete(0, tk.END)
    folder_path_entry.insert(0, selected_path)

# Function to handle folder browsing
def browse_folder():
    folder_path = filedialog.askdirectory()
    if folder_path:
        folder_path_entry.delete(0, tk.END)  # Clear the entry widget
        folder_path_entry.insert(0, folder_path)  # Insert the new folder path
        update_folder_dropdown()  # Update the dropdown with the new path

def delete_folder_path(event):
    selected_path = folder_path_dropdown.get()  # Get the currently selected path
    if selected_path:
        user_data = load_user_data()  # Load the current user data
        if selected_path in user_data["last_folders"]:
            user_data["last_folders"].remove(selected_path)  # Remove the selected path
            save_user_data(last_folders=user_data["last_folders"]) # Save the updated user data

            # If the last_folders list is now empty, remove the last_folder as well
            if not user_data["last_folders"]:
                user_data.pop("last_folder", None)  # Remove last_folder from memory
                folder_path_entry.delete(0, tk.END)  # Clear the folder path entry

            # Save the updated user data
            with open(USER_DATA_FILE, "w") as file:
                json.dump(user_data, file, indent=4)
            # Update the dropdown with the new list of folders
            folder_path_dropdown["values"] = user_data["last_folders"]
            # Clear the current selection in the dropdown
            folder_path_dropdown.set("")

            # Create a temporary popup
            popup = tk.Toplevel(root)
            popup.overrideredirect(True)  # Remove window decorations (no close button)
            popup.geometry("400x50+{}+{}".format(
                root.winfo_x() + root.winfo_width() // 2 - 50,  # Center horizontally
                root.winfo_y() + root.winfo_height() // 2 - 25  # Center vertically
            ))
            popup.configure(bg="black")  # Set background color

            # Add a label with the text "Deleted."
            label = tk.Label(popup, text="Folder Path Deleted!", fg="white", bg="black", font=("Arial", 12 ,"bold"))
            label.pack(pady=10)

            # Function to fade out the popup
            def fade_out(popup, alpha=1.0):
                if alpha <= 0:
                    popup.destroy()  # Close the popup
                    return
                popup.attributes("-alpha", alpha)  # Set transparency
                root.after(150, fade_out, popup, alpha - 0.1)  # Reduce alpha every 150ms

            # Start fading out after 500ms
            root.after(500, fade_out, popup)
  
# Function to check spelling and underline misspelled words
def check_spelling():
    user_input_box.tag_remove("misspelled", "1.0", "end")  # Clear previous misspelled tags
    text = user_input_box.get("1.0", "end-1c")  # Get the text from the input widget
    words = text.split()  # Split text into words
    start_index = "1.0"  # Start checking from the beginning

    for word in words:
        # Calculate the end index of the current word
        end_index = f"{start_index}+{len(word)}c"
        if is_english(word):
            if not spell_checker.check(word):
                user_input_box.tag_add("misspelled", start_index, end_index)  # Tag misspelled word
        else:
            pass
        
        # Move the start index to the next word
        start_index = f"{end_index}+1c"

# Function to show spelling suggestions on right-click
def show_suggestions(event):
    # Get the word under the cursor
    word_start = user_input_box.index(f"@{event.x},{event.y} wordstart")
    word_end = user_input_box.index(f"@{event.x},{event.y} wordend")
    word = user_input_box.get(word_start, word_end)

    # Check if the word is misspelled
    if not spell_checker.check(word):
        # Get suggestions for the misspelled word
        suggestions = spell_checker.suggest(word)
        
        # Create a right-click menu
        menu = Menu(root, tearoff=0)
        for suggestion in suggestions:
            menu.add_command(label=suggestion, command=lambda s=suggestion: replace_word(word_start, word_end, s))
        
        # Display the menu at the cursor position
        menu.post(event.x_root, event.y_root)

# Function to replace a misspelled word with a suggestion
def replace_word(start, end, replacement):
    user_input_box.delete(start, end)  # Delete the misspelled word
    user_input_box.insert(start, replacement)  # Insert the suggested word
    check_spelling()  # Recheck spelling after replacement

def is_english(word):
    english_chars = set("abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ")
    return all(char in english_chars for char in word)

# Copy to Clipboard button
def copy_to_clipboard():
    root.clipboard_clear()  # Clear the clipboard
    root.clipboard_append(chat_history.get("1.0", tk.END))  # Copy chat history content to clipboard

def revise_last(event):
    global do_revise
    do_revise=True
    user_input_box.delete("1.0", tk.END)  # Clear current input
    user_input_box.insert(tk.END, previous_message)  # Insert previous message

# -------------------------------------------------

# Initialize spell checker
spell_checker = enchant.Dict("en_US")
previous_message=""
do_revise=False

# Create the main window
root = tk.Tk()
root.title("AI Document Analyzer")

# Set the background color of the main window to gray
root.configure(bg="#333333")

# Configure grid weights to make the chat history box resizable
root.grid_rowconfigure(1, weight=1)  # Make row 1 (chat history) resizable
root.grid_columnconfigure(0, weight=1)  # Make column 0 resizable

# Fetch installed Ollama models
installed_models = fetch_installed_models()

# Load user data (last folder path, last selected model, and temperature)
user_data = load_user_data()
last_folder = user_data.get("last_folder", "")
last_model = user_data.get("last_model", "")
last_temperature = user_data.get("temperature", 0.7)  # Default temperature
do_mention_page = user_data.get("do_mention_page", False)  # Default checkbox state
do_read_image=True

# Dropdown for model selection
model_var = tk.StringVar(root)
model_var.set(last_model if last_model in installed_models else (installed_models[0] if installed_models else "No models found"))
model_dropdown = ttk.Combobox(root, textvariable=model_var, values=installed_models, state="readonly")
model_dropdown.grid(row=0, column=3, padx=10, pady=10, sticky="e")

# Folder path entry
folder_path_entry = ttk.Entry(root, width=50)
folder_path_entry.grid(row=0, column=0, padx=10, pady=10, sticky="ew")

# Set the last used folder path
folder_path_entry.insert(0, last_folder)

# Folder path dropdown
folder_path_dropdown = ttk.Combobox(root, values=user_data.get("last_folders", []), state="normal")
folder_path_dropdown.grid(row=0, column=0, padx=10, pady=10, sticky="ew")
folder_path_dropdown.bind("<<ComboboxSelected>>", on_folder_select)  # Bind selection event

# Set the dropdown value to the last_folder
folder_path_dropdown.set(last_folder)

# Browse button
browse_button = ttk.Button(root, text="Browse", command=browse_folder)  # Use the browse_folder function
browse_button.grid(row=0, column=1, padx=10, pady=10)

# Read button
read_button = ttk.Button(root, text="Read Documents", command=set_folder_path)
read_button.grid(row=0, column=2, padx=10, pady=10)

# Chat history display
chat_history = scrolledtext.ScrolledText(root, width=80, height=20, state=tk.DISABLED, bg="#444444", fg="white", insertbackground="white", wrap=tk.WORD)
chat_history.grid(row=1, column=0, columnspan=4, padx=10, pady=10, sticky="nsew")

# Chat history display
chat_history = scrolledtext.ScrolledText(root, width=80, height=20, state=tk.DISABLED, bg="#444444", fg="white", insertbackground="white")
chat_history.grid(row=1, column=0, columnspan=4, padx=10, pady=10, sticky="nsew")

# Sample text label
typehere_label = ttk.Label(root, text="Chat with AI here: (Shift+↵ to send | ^ to revise previous)", foreground="green",font=("Arial", 8))
typehere_label.grid(row=2, column=0, columnspan=4, padx=10, pady=(10, 0), sticky="w")  # Place between chat_history and user_input_box

# Configure tags for highlighting text
chat_history.tag_configure("fileread_tag", foreground="yellow")  # Color for file read report
chat_history.tag_configure("user_tag", foreground="cyan")  # Color for "You: "
chat_history.tag_configure("user_question", foreground="lightblue", font=("Arial", 10, "bold"))
chat_history.tag_configure("ai_tag", foreground="red")  # Color for "AI:"
chat_history.tag_configure("ai_response", foreground="white")  # Color for AI response text
chat_history.tag_configure("separator", foreground="gray")  # Color for the separator line

# User input box
user_input_box = tk.Text(root, width=60, height=3, bg="#444444", fg="white", insertbackground="white", wrap=tk.WORD)
user_input_box.grid(row=3, column=0, padx=10, pady=10, sticky="ew")

# Configure the "misspelled" tag for underlining misspelled words
user_input_box.tag_configure("misspelled", underline=True, underlinefg="lightcoral")

# Bind the spell check function to the text widget
user_input_box.bind("<KeyRelease>", lambda event: check_spelling())

# Bind the right-click event to show spelling suggestions
user_input_box.bind("<Button-3>", show_suggestions)

# Send button
send_button = ttk.Button(root, text="Ask AI", command=chat_with_ai)
send_button.grid(row=3, column=1, padx=10, pady=10)

# Clear button
clear_button = tk.Button(root, text="Clear", command=clear_chat_history ,bg="lightcoral" ,fg="black")
clear_button.grid(row=3, column=2, padx=10, pady=10)


copy_button = ttk.Button(root, text="Copy History", command=copy_to_clipboard)
copy_button.grid(row=4, column=2, padx=10, pady=10)

# Temperature slider
temperature_label = ttk.Label(root, text="Innovation Factor:", background="#333333", foreground="white")
temperature_label.grid(row=3, column=3, padx=10, pady=10)

temperature_scale = tk.Scale(root, from_=0.0, to=1.0, resolution=0.1, orient=tk.HORIZONTAL, bg="#333333", fg="white", troughcolor="#444444")
temperature_scale.set(last_temperature)  # Set the last used temperature
temperature_scale.grid(row=3, column=4, padx=10, pady=10)

# Checkbox for toggling prompt
do_mention_var = tk.BooleanVar(value=do_mention_page)  # Set the checkbox state
do_mention_checkbox = ttk.Checkbutton(
    root,
    text="Tell Which Page",
    variable=do_mention_var,
    onvalue=True,
    offvalue=False
)
do_mention_checkbox.grid(row=4, column=4, padx=10, pady=10, sticky="w")

# Bind the UP key to recall the previous user message
user_input_box.bind("<Up>", revise_last)

# Bind the Delete key to the folder_path_dropdown widget
folder_path_dropdown.bind("<Delete>", delete_folder_path)
# Bind the Enter key to trigger the chat_with_ai function
user_input_box.bind("<Shift-Return>", lambda event: chat_with_ai())

# Run the application
root.mainloop()