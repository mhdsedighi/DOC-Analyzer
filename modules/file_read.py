import os
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import win32com.client
from io import BytesIO
import base64
from modules.pdf_tools import extract_content_from_pdf
from langchain_chroma import Chroma
from langchain_community.embeddings import FastEmbedEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores.utils import filter_complex_metadata
import chromadb
import logging

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Text splitter initialized globally, configurable chunk sizes
CHUNK_SIZE = 1024  # Default to larger size from second codebase, can switch to 512
CHUNK_OVERLAP = 100 if CHUNK_SIZE == 1024 else 60
text_splitter = RecursiveCharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)

# Chroma database will be set from the main script
chroma_db = None


# Function to set the Chroma database from the main script
def set_chroma_db(db):
    global chroma_db
    chroma_db = db
    logger.info(
        f"Chroma DB set in file_read.py with persist_directory: {chroma_db._client._system.settings.persist_directory}")


# Function to check if a file is already in Chroma
def is_file_in_chroma(file_path):
    if chroma_db is None:
        raise ValueError("Chroma database not initialized")
    filename = os.path.basename(file_path)
    existing_docs = chroma_db.get(where={"filename": filename})
    logger.info(f"Checking if {filename} is in Chroma: {len(existing_docs['documents']) > 0}")
    return len(existing_docs["documents"]) > 0


# Function to extract metadata (e.g., author) from a file
def extract_author(file_path, file_ext, content=None):
    if file_ext == ".pdf":
        from pypdf import PdfReader
        try:
            reader = PdfReader(file_path)
            info = reader.metadata
            return info.get('/Author', "Unknown") if info else "Unknown"
        except Exception as e:
            logger.error(f"Error extracting author from PDF {file_path}: {e}")
            return "Unknown"
    elif file_ext == ".docx" and content:
        try:
            return content.core_properties.author or "Unknown"
        except Exception as e:
            logger.error(f"Error extracting author from DOCX {file_path}: {e}")
            return "Unknown"
    # Add more file types as needed
    return "Unknown"


# Function to extract text from a document based on its file type and store in Chroma
def extract_content_from_file(file_path, do_read_image, tesseract_path=None):
    if chroma_db is None:
        raise ValueError("Chroma database not initialized")

    filename = os.path.basename(file_path)
    file_ext = os.path.splitext(filename)[1].lower()

    # Extract content based on file type
    if file_path.endswith(".pdf"):
        text_content, image_content, word_count, file_message = extract_content_from_pdf(file_path, do_read_image,
                                                                                         tesseract_path)
    elif file_path.endswith(".docx"):
        text_content, image_content, word_count, file_message = extract_content_from_docx(file_path, do_read_image)
    elif file_path.endswith(".doc"):
        text_content, image_content, word_count, file_message = extract_content_from_doc(file_path, do_read_image)
    elif file_path.endswith(".txt"):
        text_content, image_content, word_count, file_message = extract_content_from_txt(file_path, do_read_image)
    elif file_path.endswith(".xlsx"):
        text_content, image_content, word_count, file_message = extract_content_from_xlsx(file_path, do_read_image)
    elif file_path.endswith(".xls"):
        text_content, image_content, word_count, file_message = extract_content_from_xls(file_path, do_read_image)
    elif file_path.endswith(".pptx"):
        text_content, image_content, word_count, file_message = extract_content_from_pptx(file_path, do_read_image)
    elif file_path.endswith(".ppt"):
        text_content, image_content, word_count, file_message = extract_content_from_ppt(file_path, do_read_image)
    else:
        return [], [], 0, "Unsupported File"  # Unsupported file type

    # Store extracted text in Chroma database
    if text_content:
        # Prepare texts and metadatas uniformly
        if isinstance(text_content, list) and all(isinstance(page, dict) for page in text_content):
            # PPT/PPTX case with page numbers
            texts = [page["content"] for page in text_content]
            author = extract_author(file_path, file_ext)
            metadatas = [{"filename": filename, "page": page["page"], "content": page["content"], "author": author} for page in text_content]
        else:
            # Other file types with single text or list of strings
            texts = text_content if isinstance(text_content, list) else [text_content]
            author = extract_author(file_path, file_ext)
            metadatas = [{"filename": filename, "page": i + 1, "content": text, "author": author} for i, text in enumerate(texts)]

        # Split text into chunks and filter metadata
        full_text = "\n".join(texts)
        chunks = text_splitter.split_text(full_text)
        if chunks:
            # Filter complex metadata
            chunk_docs = [dict(page_content=chunk, metadata=meta) for chunk, meta in zip(chunks, [metadatas[i % len(metadatas)] for i in range(len(chunks))])]
            filtered_chunks = filter_complex_metadata([doc for doc in chunk_docs])
            filtered_texts = [doc["page_content"] for doc in filtered_chunks]
            filtered_metadatas = [doc["metadata"] for doc in filtered_chunks]

            # Log extracted metadata for debugging
            for meta in filtered_metadatas:
                logger.info(f"Extracted metadata for {filename}: {meta}")

            logger.info(f"Adding {len(filtered_texts)} chunks for {filename} to Chroma")
            chroma_db.add_texts(texts=filtered_texts, metadatas=filtered_metadatas)

    return text_content, image_content, word_count, file_message


# Function to extract text from a DOCX file
def extract_content_from_docx(docx_path,do_read_image):
    text_content = []  # Stores extracted text
    image_content = []  # Stores extracted images (DOCX may contain embedded images)
    word_count = 0

    try:
        doc = Document(docx_path)

        # Extract all text as a single string
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        if text:
            text_content.append(text)  # Single-element list
            word_count = len(text.split())

        # Extract images if enabled
        if do_read_image:
            for rel in doc.part.rels:
                if "image" in doc.part.rels[rel].target_ref:
                    image_data = doc.part.rels[rel].target_part.blob
                    img_base64 = base64.b64encode(image_data).decode("utf-8")
                    image_content.append(img_base64)

        return text_content, image_content, word_count, ""

    except Exception as e:
        print(f"DOCX extraction failed: {e}")
        return [], [], 0, ""


# Function to extract text and images from a DOC file (old Word format)
def extract_content_from_doc(doc_path, do_read_image):
    text_content = []  # Stores extracted text
    image_content = []  # Stores extracted images
    word_count = 0

    try:
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(os.path.abspath(doc_path))

        # Extract all text as a single string
        text = doc.Content.Text.strip()
        if text:
            text_content.append(text)  # Single-element list
            word_count = len(text.split())

        # Extract images if enabled
        if do_read_image:
            for shape in doc.InlineShapes:
                if shape.Type == 3:  # Type 3 represents images
                    image = shape.PictureFormat
                    image_data = image.SaveAsPicture()  # Save the image temporarily
                    with open(image_data, "rb") as img_file:
                        img_base64 = base64.b64encode(img_file.read()).decode("utf-8")
                        image_content.append(img_base64)

        doc.Close(False)
        word.Quit()

        return text_content, image_content, word_count, ""

    except Exception as e:
        print(f"DOC extraction failed: {e}")
        return [], [], 0, ""


# Function to extract text from a TXT file
def extract_content_from_txt(txt_path, do_read_image):
    text_content = []  # Stores extracted text
    image_content = []  # No images for TXT files
    word_count = 0

    try:
        with open(txt_path, "r", encoding="utf-8") as file:
            text = file.read().strip()

        if text:
            text_content.append(text)  # Single-element list
            word_count = len(text.split())

        return text_content, image_content, word_count, ""

    except Exception as e:
        print(f"TXT extraction failed: {e}")
        return [], [], 0, ""


# Function to extract text from an XLSX file
def extract_content_from_xlsx(xlsx_path, do_read_image):
    text_content = []  # Stores extracted text
    image_content = []  # No images for XLSX files (unless explicitly handled)
    word_count = 0

    try:
        workbook = load_workbook(xlsx_path)

        # Extract all text as a single string
        sheet_text = ""
        for sheet in workbook.sheetnames:
            for row in workbook[sheet].iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    sheet_text += row_text + "\n"

        if sheet_text.strip():
            text_content.append(sheet_text)  # Single-element list
            word_count = len(sheet_text.split())

        return text_content, image_content, word_count, ""

    except Exception as e:
        print(f"XLSX extraction failed: {e}")
        return [], [], 0, ""


# Function to extract text from an XLS file (old Excel format)
def extract_content_from_xls(xls_path, do_read_image):
    text_content = []  # Stores extracted text
    image_content = []  # No images for XLS files
    word_count = 0

    try:
        excel = win32com.client.Dispatch("Excel.Application")
        excel.Visible = False
        workbook = excel.Workbooks.Open(xls_path)

        # Extract all text as a single string
        sheet_text = ""
        for sheet in workbook.Sheets:
            for row in sheet.UsedRange.Rows:
                row_text = " ".join([str(cell) for cell in row.Value if cell is not None])
                if row_text.strip():
                    sheet_text += row_text + "\n"

        if sheet_text.strip():
            text_content.append(sheet_text)  # Single-element list
            word_count = len(sheet_text.split())

        workbook.Close(False)  # Close without saving
        excel.Quit()

        return text_content, image_content, word_count, ""

    except Exception as e:
        print(f"XLS extraction failed: {e}")
        return [], [], 0, ""


# Function to extract text and images from a PPTX file
def extract_content_from_pptx(pptx_path, do_read_image):
    text_content = []  # Stores extracted text
    image_content = []  # Stores extracted images
    word_count = 0

    try:
        presentation = Presentation(pptx_path)

        for slide_num, slide in enumerate(presentation.slides, start=1):
            slide_text = ""

            # Extract text from each shape in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"

            if slide_text.strip():
                text_content.append({"page": slide_num, "content": slide_text})
                word_count += len(slide_text.split())

            # Extract images if enabled
            if do_read_image:
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        image = shape.image
                        image_stream = BytesIO(image.blob)
                        img_base64 = base64.b64encode(image_stream.getvalue()).decode("utf-8")
                        image_content.append({"page": slide_num, "content": img_base64})

        return text_content, image_content, word_count,""

    except Exception as e:
        print(f"PPTX extraction failed: {e}")
        return [], [], 0, ""


# Function to extract text and images from a PPT file (old PowerPoint format)
def extract_content_from_ppt(ppt_path, do_read_image):
    text_content = []  # Stores extracted text
    image_content = []  # Stores extracted images
    word_count = 0

    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = False
        presentation = powerpoint.Presentations.Open(ppt_path, WithWindow=False)

        for slide_num, slide in enumerate(presentation.Slides, start=1):
            slide_text = ""

            # Extract text from each shape in the slide
            for shape in slide.Shapes:
                if hasattr(shape, "TextFrame") and shape.TextFrame.HasText:
                    slide_text += shape.TextFrame.TextRange.Text + "\n"

            if slide_text.strip():
                text_content.append({"page": slide_num, "content": slide_text})
                word_count += len(slide_text.split())

            # Extract images if enabled
            if do_read_image:
                for shape in slide.Shapes:
                    if shape.Type == 13:  # Type 13 represents pictures
                        image = shape.PictureFormat
                        temp_path = os.path.join(tempfile.gettempdir(), f"ppt_image_{slide_num}.png")
                        image.SaveAsFile(temp_path)  # Save the image temporarily
                        with open(temp_path, "rb") as img_file:
                            img_base64 = base64.b64encode(img_file.read()).decode("utf-8")
                            image_content.append({"page": slide_num, "content": img_base64})
                        os.remove(temp_path)  # Clean up the temporary file

        presentation.Close()
        powerpoint.Quit()

        return text_content, image_content, word_count, ""

    except Exception as e:
        print(f"PPT extraction failed: {e}")
        return [], [], 0, ""